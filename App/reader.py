from PyPDF2 import PdfReader
from PIL import Image
from pdf2image import convert_from_path
import os
import tempfile
from pptx import Presentation
import pytesseract
import docx
# Set Tesseract executable path
pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_PATH", "tesseract")
# Set Poppler path for PDF to image conversion
poppler_path = os.getenv("POPPLER_PATH", "poppler")


def extract_doc(file):
    doc = docx.Document(file)
    text = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():  # Ignore empty paragraphs
                text.append(paragraph.text)

     # Extract text from tables
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_data:
                table_data.append(" | ".join(row_data))

        if table_data:
            text.append("\n".join(table_data))

    return "\n\n".join(text).strip()


def extract_pdf(file):
    """Extracts text from a PDF, handles multi-page PDFs, and falls back to OCR if needed."""
    reader = PdfReader(file)
    text = []

    # Check if the PDF is encrypted
    if reader.is_encrypted:
        print("❌ PDF is encrypted, skipping extraction.")
        return "❌ This PDF is encrypted and cannot be processed."


    # Extract text from all pages
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text and page_text.strip():
            text.append(page_text)
    
    # If no text was extracted, fallback to OCR
    if not text:
        print("No selectable text found. Falling back to OCR...")
        images = convert_from_path(file, poppler_path=poppler_path)
        ocr_text = " ".join([pytesseract.image_to_string(img) for img in images])
        text.append(ocr_text) # Append OCR text instead of overwriting 

    return "\n\n".join(text).strip()

def extract_ppt(file):
    #Extract text from a powerpoint file (PPT/PPTX)
    pres = Presentation(file)
    text = []

    for slide in pres.slides:
        slide_text = []

        for shape in slide.shapes :
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        # Extract Notes (if available)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_text.append(f"\n📝 Speaker Notes:\n{notes_text}")

        if slide_text:
            text.append("\n".join(slide_text))

    return "\n\n".join(text)

