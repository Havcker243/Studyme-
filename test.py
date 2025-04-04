from PyPDF2 import PdfReader
from PIL import Image
from pdf2image import convert_from_path
import pytesseract
from transformers import pipeline
import docx
from openai import OpenAI
import os 
from dotenv import load_dotenv
import hashlib
from serpapi import GoogleSearch
import redis
from pptx import Presentation
import json 
import re

import textwrap
# Set Tesseract executable path
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
poppler_path = r" C:\Program Files\poppler-24.08.0\Library\bin"
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6",framework= "pt")

load_dotenv()
api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=api_key)
serpapi_key = os.getenv("SERPAPI_KEY")
Open_router_key = os.getenv("OPEN_ROUTER_KEY")

def validate_file(file):
    """Checks if the file exists and is not empty."""
    if not os.path.exists(file):
        raise FileNotFoundError(f"Error: The file '{file}' does not exist.")
    
    if os.path.getsize(file) == 0:
        raise ValueError(f"Error: The file '{file}' is empty.")

    return True


def extract_doc(file):
    doc = docx.Document(file)
    text = []

    #Extract normal text 
    for paragraphs in doc.paragraphs:
        if paragraphs.text.strip():  # Ignore empty paragraphs
            text.append(paragraphs.text)

    # Extract text from tables
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_data:
                table_data.append(" | ".join(row_data))

        if table_data:
            text.append("\n".join(table_data))

    return "\n\n".join(text).strip()


def extract_pdf(file):
    """Extracts text from a PDF, handles multi-page PDFs, and falls back to OCR if needed."""
    reader = PdfReader(file)
    text = []

    # Extract text from all pages
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text and page_text.strip():
            text.append(page_text)
    
    # If no text was extracted, fallback to OCR
    if not text:
        print("No selectable text found. Falling back to OCR...")
        images = convert_from_path(file)
        ocr_text = " ".join([pytesseract.image_to_string(img) for img in images])
        text.append(ocr_text) # Append OCR text instea dof overwriting 

    return "\n\n".join(text).strip()

def extract_ppt(file):
    #Extract text from a powerpoint file (PPT/PPTX)
    pres = Presentation(file)
    text = []

    for slide in pres.slides:
        slide_text = []

        for shape in slide.shapes :
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        # Extract Notes (if available)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_text.append(f"\n📝 Speaker Notes:\n{notes_text}")

        if slide_text:
            text.append("\n".join(slide_text))

    return "\n\n".join(text)

def chunk_text(text, chunk_size=1024):
    """Splits text into smaller chunks with max tokens of 1024 (or less)."""
    if isinstance(text, list):
        text = "\n\n".join(text)

    words = text.split()  # Split text into words
    chunks = []
    current_chunck = []

    for word in words:
        current_chunck.append(word)

        if len(" ".join(current_chunck)) >= chunk_size:
            chunks.append(" ".join(current_chunck))
            current_chunck = []

    if current_chunck:
        chunks.append(" ".join(current_chunck))

    return chunks


def summarize_large_text(text):
    """Summarizes large text by breaking it into chunks and summarizing each separately."""
    text_chunks = chunk_text(text, chunk_size=1024)  # Break text into chunks
    summarized_text = []

    for chunk in text_chunks:
        if not chunk.strip():
            continue
        input_length = len(chunk.split())
        if input_length==0:
            continue

        max_len = min(300, int(input_length))

        summary = summarizer(chunk, max_length=max_len , min_length=int(max_len* 0.5), do_sample=False)
        summarized_text.append(summary[0]['summary_text'])

    return " ".join(summarized_text).strip()


def flashcards(text):
    response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": """ 
                         I want you to go through the whole ext and try and 
                         extract the most important things , from concepts to definitionss, examples , facts and details , think like a teacher 
                         and take in consideration mutiple test worthy question from the text and develop a a question with multi -choice and also 
                         develop a question and an answer this should be proccesed in three ways and return the output in a json format 
                         with 2 keys "MCQ" and "Cards"

                         1. **Cards**: Think in a flashcard kind of way help put the question here and then give me the answer for the question also 

                         2. **MCQ**: Should contain the question and 4 other possible answers an be different or similar but they all have to be gotten from the test 
                         the real answer will be one of the options but it should be randomized 

                         Return the result **Strictly** as a raw JSON without Markdown formatting "" like this:
                         **Output Format**
                         {
                            "Cards": [
                                    {"Question": "What is X ?", "answer": "X is ....."},
                                    {"Question": "What is X ?", "answer": "X is ....."}, 
                                    {"Question": "What is X ?", "answer": "X is ....."}, 
                                    {"Question": "What is X ?", "answer": "X is ....."} 
                                    ], 

                            "MCQ": [
                                    {"Question": "Questions Here", "options":[Option1, Option2, Option3 , Option4], "correct_answer": "Option3"}, 
                                    {"Question": "Questions Here", "options":[Option1, Option2, Option3 , Option4], "correct_answer": "Option2"}, 
                                    {"Question": "Questions Here", "options":[Option1, Option2, Option3 , Option4], "correct_answer": "Option1"}, 
                                    {"Question": "Questions Here", "options":[Option1, Option2, Option3 , Option4], "correct_answer": "Option4"}, 
                                    {"Question": "Questions Here", "options":[Option1, Option2, Option3 , Option4], "correct_answer": "Option3"},
                                    {"Question": "Questions Here", "options":[Option1, Option2, Option3 , Option4], "correct_answer": "Option1"}
                                    
                                    ]
                         } 
                         """
                    },

                    {
                        "role": "user",
                        "content": text
                    }
                ]
        )
    answer = response.choices[0].message.content.strip()
    try:
        clean_text = answer.replace("```json", "").replace("```", "").strip()
        parsed_reponse = json.loads(clean_text)
    except json.JSONDecodeError:
        print("Error: Deepseek did not return valid JSON")
        print("Raw Response:", answer)
        return None
    
    cards = parsed_reponse.get("Cards", [])
    MCQ = parsed_reponse.get("MCQ", [])
    
    return {"Cards": cards, "MCQ": MCQ}






def explain(text):
    response = client .chat.completions.create(
        model= "gpt-4o",
        messages=[
            {
              "role":"developer",
              "content":"""
                I want you to process this text in three ways and return the output in JSON format with three keys: 
                'bullets' , "Notes "and 'explanation'. 

                1. **Bullets**: Extract the key terms from this text in a **simple bullet format**. 
                 - Only include the **most relevant** key terms, important concepts, and main points. 
                 - Do NOT explain—just list the terms.

                2. **Explanation**: Take the same text and provide a **detailed** breakdown.
                 - Explain all the **main points**, important details, and core ideas from the text, I want you to touch on it each select part 
                - Use two to three small and larges **examples** to clarify meaning of each part also for easy understanding 
                 - Ensure the explanation is **easy enough for a 5-year-old**, but also **detailed enough for a very genus and gifted college student**.
                 - Always **expand on key ideas** instead of summarizing briefly.

                3. **Notes** : I want you to give expressive notes to try and take the most important concepts and points 
                  - Let it be be optimized , in a list , and still talk about each concepts 
                  - Check and make sure that is is clean and understanable 
                  - Work and try to thinkin the way a teache to try and set test and exams from thuis text when it comes to note taking , be like a harvard student who went to yale and MIT for post grdaduates 
        
                Return the result  **strictly** as a raw JSON without Markdown formatting "" like this :
                {{
                    "bullets": ["key term 1", "key term 2", "key term 3"],
                    "explanation": "Full detailed explanation here.",
                    "Notes" : "Fulle detailed explanation here. "
                }}
                """
            },
            {
                "role":"user",
                "content": text
            }
        ]
    )
    answer = response.choices[0].message.content.strip()
    try:

        # Attempt to parse as JSON
        parsed_response = json.loads(answer)
        return parsed_response
    except json.JSONDecodeError:
        print("Error: OpenAI did not return valid JSON. Full response:")
        print(answer)
        return None



def search_using_bullets(parsed_response):
    """
    Takes the parsed OpenAI response, extracts all the 'bullets' list,
    and performs a web search for each term using SerpAPI.
    """
    if not parsed_response or "bullets" not in parsed_response:
        return {"error": "No valid bullets extracted from OpenAI response."}

    bullets = parsed_response["bullets"]  # Extract key terms
    search_results = {}

    serpapi_key = os.getenv("SERPAPI_KEY")  # Ensure API key is set

    for term in bullets:
        params = {
            "q": term,
            "api_key": serpapi_key,
            "num": 3  # Limit to top 6 search results per term
        }
        search = GoogleSearch(params)
        results = search.get_dict().get("organic_results", [])
        search_results[term] = [{"title": result["title"], "link": result["redirect_link"]} for result in results] # Store search results per key term

    return search_results

    



# Load PDF
file = "./TestDocs/p.pdf"
text = extract_pdf(file)
print(text)
print("\n")
print("\n")
print("\n")
t = chunk_text(text, 1024)
print("Here are the chuncks\n")
print(t ,"\n")
# q , a = "", ""
# print("\n" + "-" * 50 + "\n")
# print("Here are the Flashcards\n")
# j = flashcards(t)
# c2 = j['Cards']
# m2 = j['MCQ']

# for mp in c2:
#     q , a = mp['Question'], mp['answer']
#     print("Question:", q, "\nAnswer:", a)

# for np in m2:
#     q , o = np['Question'], np['options']

#     print("Question:", q)
#     for n in range(len(o)):
#         print("\nOptions:", o[n])



# print("\n")
# print("\n")
# print("\n")
# print("\n")
# print("\n")
# print("\n")
# print("\n" + "-" * 50 + "\n")

h = summarize_large_text(t)
print("Here are the summarized\n")
print(h)
# Summarize in chunks
final_summary = h
explained = explain(text)
search_results = search_using_bullets(explained)
print(explained["bullets"])
print(explained["explanation"])
print( "\n" +  "\n")
print(search_results)
for term, results in search_results.items():
    print(f"Here are results based on this term: {term}")  # Print the search term
    print("The links are below:")
    
    for result in results:
        print(result["title"])
        print("\n")
        print(result["link"])
    
    print("\n" + "-" * 50 + "\n")